import streamlit as st
import PyPDF2
import docx  
import pptx
import pandas as pd
import google.generativeai as genai
import nltk
from nltk.tokenize import sent_tokenize
import os
from dotenv import load_dotenv
import time
from google.api_core.exceptions import ResourceExhausted
import json  # <-- Add this import

nltk.download('punkt')
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=GOOGLE_API_KEY)

# Add dropdown for Gemini model selection
model_options = [
    "gemini-1.5-flash",
    "gemini-1.5-flash-latest",
    "gemini-1.5-pro",
    "gemini-1.5-pro-latest",
    "gemini-2.5-flash"
]
selected_model = st.sidebar.selectbox("Select Gemini Model", model_options, index=model_options.index("gemini-2.5-flash"))
model = genai.GenerativeModel(selected_model)

def extract_text(file):
    if file.type == "application/pdf":
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = pptx.Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        df = pd.read_excel(file)
        text = df.to_string()
        return text
    elif file.type == "text/plain":
        return file.read().decode("utf-8")
    elif file.type == "application/json":
        try:
            file_content = file.read().decode("utf-8")
            data = json.loads(file_content)
            # Convert JSON to a readable string (pretty print)
            text = json.dumps(data, indent=2)
            return text
        except Exception as e:
            return f"Error reading JSON: {e}"
    else:
        return ""

def summarize_document(text):
    prompt = f"Summarize the following document in no more than 150 words:\n{text}"
    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        st.error(f"Gemini error: {e}")
        return "Unable to summarize. Please check your API key, quota, or try again later."

def answer_question(document, question, history=[]):
    prompt = f"""You are a research assistant. Answer the following question based only on the provided document. 
Document: {document}
Question: {question}
If you answer, include a brief justification referencing the document (e.g., 'This is supported by paragraph 3 of section 1...') and highlight the supporting snippet.
"""
    if history:
        prompt += f"\nPrevious interactions:\n" + "\n".join(history)
    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        st.error(f"Gemini error: {e}")
        return "Unable to answer. Please check your API key, quota, or try again later."

def generate_logic_questions(document):
    prompt = f"""Read the following document and generate three logic-based or comprehension-focused questions that require reasoning. 
Document: {document}
For each question, provide only the question text.
"""
    try:
        response = model.generate_content(prompt)
        questions = [q.strip("- ").strip() for q in response.text.strip().split('\n') if q.strip()]
        return questions[:3]
    except Exception as e:
        st.error(f"Gemini error: {e}")
        return ["Unable to generate questions. Please check your API key, quota, or try again later."] * 3

def evaluate_answer(document, question, user_answer):
    prompt = f"""Given the document and the user's answer to the question, evaluate the answer for correctness and provide feedback. 
Document: {document}
Question: {question}
User's Answer: {user_answer}
Provide feedback with justification referencing the document and highlight the supporting snippet.
"""
    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        st.error(f"Gemini error: {e}")
        return "Unable to evaluate. Please check your API key, quota, or try again later."

st.set_page_config(page_title="📊 CodHelp Smart Assistant for Research Summarization", layout="wide")
st.title("📊 CodHelp Smart Assistant for Research Summarization")
st.markdown("<div style='text-align: left; font-size: 16px; color: #888;'>Built by Vikas Singh</div>", unsafe_allow_html=True)

if "file_type_stats" not in st.session_state:
    st.session_state.file_type_stats = {}
if "total_questions" not in st.session_state:
    st.session_state.total_questions = {}
if "generation_times" not in st.session_state:
    st.session_state.generation_times = {}

st.sidebar.header("Upload Document")
uploaded_file = st.sidebar.file_uploader(
    "Choose a PDF, Word, PPT, Excel, TXT, or JSON file",
    type=["pdf", "docx", "pptx", "xlsx", "txt", "json"]
)

# Add this at the top after imports to handle Streamlit mobile network errors gracefully
st.set_option('client.showErrorDetails', True)


# Hosted dashboard link (opens in new tab)
st.sidebar.markdown(
    "<center>"
    "<div style='margin-top:8px; margin-bottom:8px;'>"
    "<a href='https://codhelassistantdashboard.streamlit.app/' target='_blank' rel='noopener'>"
    "<button style='background-color:#0e6fff;color:#fff;padding:8px 12px;border-radius:6px;border:none;cursor:pointer;'>Open Hosted Dashboard</button>"
    "</a>"
    "</div>"
    "</center>",
    unsafe_allow_html=True
)

# Instruct users about mobile network issues
st.sidebar.markdown("**Note:** If you see a 'Network Error' (AxiosError) on mobile, please check your internet connection and ensure your device can reach external APIs. Try switching to a stable WiFi or desktop browser for best results.")

if uploaded_file:
    ext = uploaded_file.name.split('.')[-1].lower()
    st.session_state.file_type_stats[ext] = st.session_state.file_type_stats.get(ext, 0) + 1

    document_text = extract_text(uploaded_file)
    st.subheader("Auto Summary (≤150 words)")
    with st.spinner("Summarizing document..."):
        summary = summarize_document(document_text)
    st.info(summary)
    
    mode = st.radio("Choose Interaction Mode", ["Ask Anything", "Challenge Me"])
    if "history" not in st.session_state:
        st.session_state.history = []
    
    if ext not in st.session_state.total_questions:
        st.session_state.total_questions[ext] = 0
    if ext not in st.session_state.generation_times:
        st.session_state.generation_times[ext] = []

    if mode == "Ask Anything":
        st.subheader("Ask Anything")
        user_question = st.text_input("Your question about the document:")
        if user_question:
            start_time = time.time()
            with st.spinner("Generating answer..."):
                answer = answer_question(document_text, user_question, st.session_state.history)
            elapsed = time.time() - start_time
            st.session_state.generation_times[ext].append(elapsed)
            st.markdown("**Answer:**")
            st.write(answer)
            st.session_state.history.append(f"Q: {user_question}\nA: {answer}")
            st.session_state.total_questions[ext] += 1

    elif mode == "Challenge Me":
        st.subheader("Challenge Me: Logic-Based Questions")
        st.info("Note: This feature may quickly exhaust your Gemini API quota. If you hit quota errors, try switching to a different Gemini model or wait for quota reset.")
        if "logic_questions" not in st.session_state or st.button("Generate New Questions"):
            with st.spinner("Generating questions..."):
                st.session_state.logic_questions = generate_logic_questions(document_text)
            st.session_state.user_answers = [""] * 3
            st.session_state.feedbacks = [""] * 3

        for idx, q in enumerate(st.session_state.logic_questions):
            st.markdown(f"**Question {idx+1}:** {q}")
            user_ans = st.text_input(f"Your answer to Q{idx+1}:", key=f"ans_{idx}")
            if user_ans:
                start_time = time.time()
                with st.spinner("Evaluating answer..."):
                    feedback = evaluate_answer(document_text, q, user_ans)
                elapsed = time.time() - start_time
                st.session_state.generation_times[ext].append(elapsed)
                st.markdown("**Feedback:**")
                st.write(feedback)
                st.session_state.user_answers[idx] = user_ans
                st.session_state.feedbacks[idx] = feedback
                st.session_state.total_questions[ext] += 1

# Show only the data table with file type, total questions, and avg. generation time
if st.session_state.file_type_stats:
    st.subheader("Assistant Data Table")
    file_types = list(st.session_state.file_type_stats.keys())
    total_questions = [st.session_state.total_questions.get(ft, 0) for ft in file_types]
    avg_gen_time = [
        round(sum(st.session_state.generation_times.get(ft, [])) / len(st.session_state.generation_times.get(ft, [])), 2)
        if st.session_state.generation_times.get(ft, []) else 0
        for ft in file_types
    ]

    st.dataframe({
        "File Type": file_types,
        "Total Questions": total_questions,
        "Avg. Generation Time (s)": avg_gen_time
    })

st.sidebar.markdown("---")
st.sidebar.markdown("**Instructions:**")
st.sidebar.markdown("""
1. Upload a PDF, Word, PPT, Excel, TXT, or JSON document.
2. View the auto summary.
3. Choose 'Ask Anything' to ask questions about the document.
4. Choose 'Challenge Me' to answer logic-based questions and get feedback.

To run the CodHelp Smart Assistant software:

1. Open a terminal and navigate to your project directory:
   cd g:\react\aiProjectinternship

2. Run the Streamlit app:
   streamlit run SmartAssistantUI.py

3. A browser window will open with the assistant UI.  
   - Use the sidebar to upload a document.
   - View the auto summary.
   - Choose "Ask Anything" or "Challenge Me" to interact with the assistant.

Make sure your `.env` file contains your GOOGLE_API_KEY.

1. Open a terminal and navigate to your project directory:
   cd g:\react\aiProjectinternship

2. Run the Streamlit app:
   streamlit run SmartAssistantUI.py

3. A browser window will open with the assistant UI.  
   - Use the sidebar to upload a document.
   - View the auto summary.
   - Choose "Ask Anything" or "Challenge Me" to interact with the assistant.

Make sure your `.env` file contains your GOOGLE_API_KEY.

   - Choose "Ask Anything" or "Challenge Me" to interact with the assistant.

Make sure your `.env` file contains your GOOGLE_API_KEY.
""")
st.sidebar.markdown("**Instructions:**")
st.sidebar.markdown("""
1. Upload a PDF, Word, PPT, Excel, TXT, or JSON document.
2. View the auto summary.
3. Choose 'Ask Anything' to ask questions about the document.
4. Choose 'Challenge Me' to answer logic-based questions and get feedback.

To run the CodHelp Smart Assistant software:

1. Open a terminal and navigate to your project directory:
   cd g:\react\aiProjectinternship

2. Run the Streamlit app:
   streamlit run SmartAssistantUI.py

3. A browser window will open with the assistant UI.  
   - Use the sidebar to upload a document.
   - View the auto summary.
   - Choose "Ask Anything" or "Challenge Me" to interact with the assistant.

Make sure your `.env` file contains your GOOGLE_API_KEY.

1. Open a terminal and navigate to your project directory:
   cd g:\react\aiProjectinternship

2. Run the Streamlit app:
   streamlit run SmartAssistantUI.py

3. A browser window will open with the assistant UI.  
   - Use the sidebar to upload a document.
   - View the auto summary.
   - Choose "Ask Anything" or "Challenge Me" to interact with the assistant.

Make sure your `.env` file contains your GOOGLE_API_KEY.

   - Choose "Ask Anything" or "Challenge Me" to interact with the assistant.

Make sure your `.env` file contains your GOOGLE_API_KEY.
""")
